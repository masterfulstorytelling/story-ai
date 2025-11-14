"""File parser for PDF, PPTX, and DOCX files."""

import re
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

from src.models.parsed_content import ParsedContent
from src.models.page import Page
from src.models.section import Section
from src.utils.logger import get_logger

logger = get_logger(__name__)


class UnsupportedFileFormatError(Exception):
    """Exception raised for unsupported file formats."""

    pass


class FileParsingError(Exception):
    """Exception raised for file parsing errors."""

    pass


def count_words(text: str) -> int:
    """Count words in text."""
    return len(re.findall(r"\b\w+\b", text))


def parse_pdf(file_content: bytes, filename: str) -> ParsedContent:
    """Parse PDF file content."""
    try:
        # Open PDF from bytes
        pdf_doc = fitz.open(stream=file_content, filetype="pdf")

        text_parts = []
        pages = []

        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            page_text = page.get_text()
            text_parts.append(page_text)

            pages.append(
                Page(
                    page_number=page_num + 1,
                    text=page_text,
                    word_count=count_words(page_text),
                )
            )

        full_text = "\n\n".join(text_parts)
        total_words = count_words(full_text)

        pdf_doc.close()

        return ParsedContent(
            filename=filename,
            file_type="pdf",
            text=full_text,
            word_count=total_words,
            pages=pages,
        )
    except Exception as e:
        raise FileParsingError(f"Error parsing PDF {filename}: {str(e)}")


def parse_pptx(file_content: bytes, filename: str) -> ParsedContent:
    """Parse PPTX file content."""
    try:
        from io import BytesIO

        # Open PPTX from bytes
        pptx_file = BytesIO(file_content)
        presentation = Presentation(pptx_file)

        text_parts = []
        sections = []

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text_parts.append(shape.text.strip())

            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                text_parts.append(slide_text)
                sections.append(
                    Section(
                        title=f"Slide {slide_num}",
                        text=slide_text,
                        page_number=slide_num,
                        order=slide_num,
                    )
                )

        full_text = "\n\n".join(text_parts)
        total_words = count_words(full_text)

        return ParsedContent(
            filename=filename,
            file_type="pptx",
            text=full_text,
            word_count=total_words,
            sections=sections,
        )
    except Exception as e:
        raise FileParsingError(f"Error parsing PPTX {filename}: {str(e)}")


def parse_docx(file_content: bytes, filename: str) -> ParsedContent:
    """Parse DOCX file content."""
    try:
        from io import BytesIO

        # Open DOCX from bytes
        docx_file = BytesIO(file_content)
        document = Document(docx_file)

        text_parts = []
        sections = []

        for para_num, paragraph in enumerate(document.paragraphs, 1):
            para_text = paragraph.text.strip()
            if para_text:
                text_parts.append(para_text)
                # Use paragraph style as section title if available
                section_title = None
                if paragraph.style and paragraph.style.name.startswith("Heading"):
                    section_title = para_text[:100]  # Use first part as title

                sections.append(
                    Section(
                        title=section_title,
                        text=para_text,
                        order=para_num,
                    )
                )

        full_text = "\n\n".join(text_parts)
        total_words = count_words(full_text)

        return ParsedContent(
            filename=filename,
            file_type="docx",
            text=full_text,
            word_count=total_words,
            sections=sections,
        )
    except Exception as e:
        raise FileParsingError(f"Error parsing DOCX {filename}: {str(e)}")


def parse_file(file_content: bytes, filename: str) -> ParsedContent:
    """
    Parse file content based on file extension.

    Args:
        file_content: File content as bytes
        filename: Name of the file

    Returns:
        ParsedContent object

    Raises:
        UnsupportedFileFormatError: If file format is not supported
        FileParsingError: If parsing fails
    """
    filename_lower = filename.lower()

    if filename_lower.endswith(".pdf"):
        return parse_pdf(file_content, filename)
    elif filename_lower.endswith(".pptx"):
        return parse_pptx(file_content, filename)
    elif filename_lower.endswith(".docx"):
        return parse_docx(file_content, filename)
    else:
        raise UnsupportedFileFormatError(
            f"Unsupported file format: {filename}. Supported formats: PDF, PPTX, DOCX"
        )
